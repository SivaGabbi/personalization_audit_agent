import os
import time
import uuid
import threading
import json
import pandas as pd
from urllib.parse import urlparse
from flask import Flask, request, jsonify
from flask_cors import CORS
from selenium import webdriver
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
from skimage.metrics import structural_similarity as ssim
import numpy as np
from scipy.spatial.distance import cosine
from PIL import Image, ImageChops
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt

# --- In-memory storage for tasks and knowledge base ---
tasks = {}
kb_df = None

# --- Configuration ---
VENDORS_FILE = 'vendors.csv'
KB_FILE = 'personalization_kb.csv'
RECOMMENDATION_KEYWORDS = [
    'recommended for you', 'inspired by your browsing', 'recently viewed',
    'customers also bought', 'frequently bought together', 'similar products',
    'you might also like', 'because you viewed', 'complete your purchase'
]

# --- RAG/LLM Setup ---
llm_model = None
embedding_model = None
try:
    api_key = os.environ.get("GEMINI_API_KEY")
    if not api_key:
        raise ValueError("GEMINI_API_KEY environment variable not set.")
    genai.configure(api_key=api_key)
    llm_model = genai.GenerativeModel('gemini-1.5-flash-latest')
    embedding_model = 'models/embedding-001'
except Exception as e:
    print(f"⚠️ RAG/LLM Warning: {e}")

# --- Helper Functions ---
def initialize_knowledge_base():
    global kb_df
    try:
        kb_df = pd.read_csv(KB_FILE)
        if embedding_model:
            print("🧠 Embedding Knowledge Base...")
            kb_df['combined_text'] = kb_df['use_case'] + ": " + kb_df['best_practice']
            embeddings = genai.embed_content(model=embedding_model, content=kb_df['combined_text'].tolist(), task_type="retrieval_document")
            kb_df['embedding'] = embeddings['embedding']
            print("   - ✅ Knowledge Base embedded successfully.")
    except Exception as e:
        print(f"❌ Failed to initialize Knowledge Base: {e}")
        kb_df = None

def update_task(task_id, status=None, message=None, progress=None, result=None):
    if task_id in tasks:
        if status: tasks[task_id]['status'] = status
        if message:
            timestamp = time.strftime('%H:%M:%S')
            tasks[task_id]['logs'].append(f"[{timestamp}] {message}")
        if progress is not None: tasks[task_id]['progress'] = progress
        if result: tasks[task_id]['result'] = result

def retrieve_relevant_context(findings_summary_text, top_k=3):
    if kb_df is None or 'embedding' not in kb_df.columns or not embedding_model:
        return []
    try:
        findings_embedding = genai.embed_content(model=embedding_model, content=findings_summary_text, task_type="retrieval_query")['embedding']
        kb_df['similarity'] = kb_df['embedding'].apply(lambda emb: 1 - cosine(emb, findings_embedding))
        top_k_df = kb_df.nlargest(top_k, 'similarity')
        context = []
        for _, row in top_k_df.iterrows():
            context.append({"use_case": row['use_case'], "best_practice": row['best_practice']})
        return context
    except Exception as e:
        print(f"Error during context retrieval: {e}")
        return []

def generate_recommendations_with_llm(task_id, score_report):
    update_task(task_id, message="[Agent 2] Analyzing score report with RAG engine...")
    if not llm_model:
        return [{"recommendation": "LLM model not available.", "justification": "Could not generate recommendations."}]

    findings_summary = []
    for category in score_report['breakdown'].values():
        for item in category:
            if item['score'] > 0:
                findings_summary.append(f"{item['name']} (Score: {item['score']:.0f}/{item['max_score']})")
    findings_summary_text = ", ".join(findings_summary)
    
    rag_context = retrieve_relevant_context(findings_summary_text)
    update_task(task_id, message=f"[Agent 2] Retrieved {len(rag_context)} relevant context(s) from knowledge base.")

    prompt = f"""
    You are Agent 2, a world-class Personalization Advisor.
    Agent 1 conducted a technical analysis and produced the following score report:
    - Overall Score: {score_report['total_score']:.0f}/100
    - Key Findings: {findings_summary_text}

    **Your Task:**
    You MUST act as an expert consultant and generate recommendations based **only** on the following retrieved knowledge base articles. For each recommendation, you must explicitly reference the 'use_case' it relates to.

    **Retrieved Knowledge Base Articles:**
    {json.dumps(rag_context, indent=2)}

    Generate a JSON object with a key "recommendations". This should be a list of **10 to 15** concise, actionable recommendations where each item is an OBJECT containing two keys:
    1. "recommendation" (a string with the actionable advice derived from a 'best_practice')
    2. "justification" (a string explaining how this advice addresses a key finding, referencing the corresponding 'use_case')
    """
    try:
        response = llm_model.generate_content(prompt)
        cleaned_response = response.text.strip().replace('```json', '').replace('```', '')
        llm_result = json.loads(cleaned_response)
        update_task(task_id, message="   - ✅ Recommendations generated successfully.")
        return llm_result.get('recommendations', [])
    except Exception as e:
        update_task(task_id, message=f"   - ❌ LLM Error: {e}")
        return [{"recommendation": "An error occurred while generating recommendations.", "justification": str(e)}]

def create_and_save_presentation(result_data, filename):
    prs = Presentation()
    score_report = result_data.get('score_report', {})
    recommendations = result_data.get('recommendations', [])
    
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title, subtitle = slide.shapes.title, slide.placeholders[1]
    url = score_report.get("url", "Website")
    title.text = "Personalization & Performance Audit"
    subtitle.text = f"An analysis of {url}\nGenerated on {pd.Timestamp.now().strftime('%Y-%m-%d')}"
    
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_shape.text = "Executive Summary: Score Breakdown"
    title_shape.text_frame.paragraphs[0].font.bold = True; title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    tf = txBox.text_frame
    tf.text = f"Overall Score: {score_report.get('total_score', 0):.0f}/100"
    tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.size = Pt(24)
    
    for category_name, items in score_report.get('breakdown', {}).items():
        p = tf.add_paragraph()
        p.text = f"\n{category_name.capitalize()} ({score_report.get(f'{category_name}_score', 0):.0f} pts)"
        p.font.bold = True
        for item in items:
            p = tf.add_paragraph()
            p.text = f"{item['name']}: {item['score']:.0f}/{item['max_score']} - {item.get('details', '')}"
            p.level = 1
            
    for i, rec_obj in enumerate(recommendations, 1):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Recommendation #{i}"
        tf = slide.placeholders[1].text_frame
        tf.clear()
        p = tf.add_paragraph()
        p.text = rec_obj.get('recommendation', str(rec_obj))
        p.font.bold = True; p.font.size = Pt(20)
        p = tf.add_paragraph()
        p.text = rec_obj.get('justification', '')
        p.font.size = Pt(16)
        
    if not os.path.exists('reports'):
        os.makedirs('reports')
    prs.save(filename)

def detect_vendors(driver, task_id):
    update_task(task_id, message="[Agent 1] Scanning for A/B testing and personalization tools...")
    try: vendors_df = pd.read_csv(VENDORS_FILE)
    except FileNotFoundError:
        update_task(task_id, message=f"   - ⚠️  Warning: '{VENDORS_FILE}' not found. Skipping tool detection.")
        return []
    found_vendors = []
    scripts = driver.execute_script("return Array.from(document.scripts).map(s => s.src);")
    script_sources_text = " ".join(filter(None, scripts))
    for _, vendor in vendors_df.iterrows():
        for signature in vendor['signatures'].split(','):
            is_present = False
            if signature in script_sources_text: is_present = True
            else:
                try:
                    if driver.execute_script(f"return typeof {signature} !== 'undefined'"): is_present = True
                except: pass
            if is_present:
                vendor_name = f"{vendor['vendor_name']} ({vendor['category']})"
                if vendor_name not in found_vendors: found_vendors.append(vendor_name)
    if found_vendors: update_task(task_id, message=f"   - ✅ Found tools: {', '.join(found_vendors)}")
    else: update_task(task_id, message="   - ⚪ No specific vendor tools detected.")
    return found_vendors

def check_for_recommendations(driver, task_id):
    try:
        body_text = driver.find_element(By.TAG_NAME, 'body').text.lower()
        for keyword in RECOMMENDATION_KEYWORDS:
            if keyword in body_text:
                update_task(task_id, message=f"   - ✅ Found recommendation keyword: '{keyword}'")
                return True
    except Exception: return False
    return False

def find_product_links(driver, base_url, task_id, limit=3):
    update_task(task_id, message="[Agent 1] Searching for product links...")
    link_patterns = ['/products/', '/p/', '/dp/', '/item/']
    links = []
    WebDriverWait(driver, 15).until(EC.presence_of_element_located((By.TAG_NAME, "a")))
    all_links = driver.find_elements(By.TAG_NAME, 'a')
    for a in all_links:
        href = a.get_attribute('href')
        if href and any(pattern in href for pattern in link_patterns):
            if href.startswith('http') and urlparse(href).netloc == urlparse(base_url).netloc:
                if href not in links: links.append(href)
    if not links:
        update_task(task_id, message="   - No common product patterns found. Falling back to any valid link.")
        for a in all_links:
            href = a.get_attribute('href')
            if href and href.startswith('http') and urlparse(href).netloc == urlparse(base_url).netloc and href != base_url:
                if href not in links: links.append(href)
    update_task(task_id, message=f"   - Found {len(links)} potential links.")
    return links[:limit]

def compare_images(path_before, path_after, diff_path, task_id):
    try:
        before_img = Image.open(path_before).convert('RGB')
        width, height = before_img.size
        if width < 7 or height < 7:
            update_task(task_id, message=f"   - ⚠️  Warning: Screenshot is too small. Skipping comparison.")
            return 0.0
        after_img = Image.open(path_after).convert('RGB')
        if before_img.size != after_img.size: after_img = after_img.resize(before_img.size)
        diff_img = ImageChops.difference(before_img, after_img)
        diff_img.save(diff_path)
        before_arr, after_arr = np.array(before_img), np.array(after_img)
        similarity_score = ssim(before_arr, after_arr, channel_axis=-1, data_range=before_arr.max() - before_arr.min())
        return (1 - similarity_score) * 100
    except FileNotFoundError: return 0

def get_performance_score(driver, task_id):
    update_task(task_id, message="[Agent 1] Measuring site performance...")
    try:
        timing_js = "return window.performance.getEntriesByType('navigation')[0];"
        nav_timing = driver.execute_script(timing_js)
        fcp = nav_timing.get('firstContentfulPaint', 3000)
        dom_interactive = nav_timing.get('domInteractive', 3000)
        update_task(task_id, message=f"   - First Contentful Paint (FCP): {fcp:.0f}ms")
        update_task(task_id, message=f"   - DOM Interactive: {dom_interactive:.0f}ms")
        fcp_score = 10 if fcp < 1800 else (5 if fcp < 3000 else 0)
        dom_interactive_score = 10 if dom_interactive < 2000 else (5 if dom_interactive < 4000 else 0)
        return {"fcp_score": fcp_score, "dom_interactive_score": dom_interactive_score, "fcp_value": fcp, "dom_interactive_value": dom_interactive}
    except Exception as e:
        update_task(task_id, message=f"   - ⚠️ Could not retrieve performance metrics: {e}")
        return {"fcp_score": 0, "dom_interactive_score": 0, "fcp_value": 0, "dom_interactive_value": 0}

def get_visual_clarity_score(screenshot_path, task_id):
    update_task(task_id, message="[Agent 1] Analyzing visual clarity...")
    try:
        img = Image.open(screenshot_path).convert('RGB')
        img = img.resize((400, int(400 * img.height / img.width)))
        colors = img.getcolors(img.width * img.height)
        num_dominant_colors = len([c for c in colors if c[0] > (img.width * 0.01)])
        update_task(task_id, message=f"   - Found {num_dominant_colors} dominant colors.")
        color_score = 10 if num_dominant_colors <= 10 else (5 if num_dominant_colors <= 20 else 0)
        pixels = np.array(img) / 255.0
        luminance = 0.299 * pixels[:,:,0] + 0.587 * pixels[:,:,1] + 0.114 * pixels[:,:,2]
        dark_pixels, light_pixels = np.sum(luminance < 0.2), np.sum(luminance > 0.8)
        total_pixels = luminance.size
        contrast_ratio = min(dark_pixels, light_pixels) / total_pixels
        update_task(task_id, message=f"   - Contrast distribution ratio: {contrast_ratio:.2f}")
        contrast_score = 5 if contrast_ratio > 0.1 else (2 if contrast_ratio > 0.05 else 0) # Max 5 for this simple heuristic
        return {"contrast_score": contrast_score, "color_score": color_score, "num_colors": num_dominant_colors, "contrast_ratio": contrast_ratio}
    except Exception as e:
        update_task(task_id, message=f"   - ⚠️ Could not analyze visual clarity: {e}")
        return {"contrast_score": 0, "color_score": 0, "num_colors": 0, "contrast_ratio": 0}

def attempt_add_to_cart(driver, task_id):
    update_task(task_id, message="[Agent 1] Attempting to add product to cart...")
    selectors = ["button[data-testid='add-to-cart-button']", "button[id*='add-to-cart']", "button[class*='add-to-cart']", "button[name='add']", "//button[contains(translate(., 'ABCDEFGHIJKLMNOPQRSTUVWXYZ', 'abcdefghijklmnopqrstuvwxyz'), 'add to cart')]", "//button[contains(translate(., 'ABCDEFGHIJKLMNOPQRSTUVWXYZ', 'abcdefghijklmnopqrstuvwxyz'), 'add to bag')]", "//button[contains(translate(., 'ABCDEFGHIJKLMNOPQRSTUVWXYZ', 'abcdefghijklmnopqrstuvwxyz'), 'add to basket')]"]
    cart_selectors = ["[data-testid='cart-count']", "[class*='cart-count']", "[id*='cart-count']", "[data-cy='cart-count']"]
    try:
        initial_cart_count = 0
        for selector in cart_selectors:
            try:
                count_element = driver.find_element(By.CSS_SELECTOR, selector)
                initial_cart_count = int(count_element.text)
                break
            except: continue
        for selector in selectors:
            try:
                add_to_cart_button = driver.find_element(By.XPATH, selector) if selector.startswith("//") else driver.find_element(By.CSS_SELECTOR, selector)
                add_to_cart_button.click()
                update_task(task_id, message="   - ✅ Clicked 'Add to Cart' button.")
                time.sleep(3)
                body_text = driver.find_element(By.TAG_NAME, 'body').text.lower()
                if "added to cart" in body_text or "added to bag" in body_text or "added to your basket" in body_text:
                    update_task(task_id, message="   - ✅ Verified: Success message found.")
                    return True
                for cart_selector in cart_selectors:
                    try:
                        count_element = driver.find_element(By.CSS_SELECTOR, cart_selector)
                        new_cart_count = int(count_element.text)
                        if new_cart_count > initial_cart_count:
                            update_task(task_id, message="   - ✅ Verified: Cart item count increased.")
                            return True
                    except: continue
                update_task(task_id, message="   - ⚪ Clicked button, but could not verify success.")
                return False
            except: continue
        update_task(task_id, message="   - ❌ Could not find an 'Add to Cart' button.")
        return False
    except Exception as e:
        update_task(task_id, message=f"   - ❌ Error during add to cart attempt: {e}")
        return False

def attempt_newsletter_signup(driver, task_id):
    update_task(task_id, message="[Agent 1] Attempting newsletter signup...")
    try:
        email_input = driver.find_element(By.CSS_SELECTOR, "input[type='email'], input[name='email']")
        submit_button = email_input.find_element(By.XPATH, "./following::button[contains(translate(., 'ABCDEFGHIJKLMNOPQRSTUVWXYZ', 'abcdefghijklmnopqrstuvwxyz'), 'sign up')] | ./following::button[contains(translate(., 'ABCDEFGHIJKLMNOPQRSTUVWXYZ', 'abcdefghijklmnopqrstuvwxyz'), 'subscribe')] | ./following::input[@type='submit']")
        fake_email = f"testuser_{uuid.uuid4().hex[:8]}@example.com"
        email_input.send_keys(fake_email)
        submit_button.click()
        update_task(task_id, message=f"   - ✅ Found form and submitted email.")
        time.sleep(3)
        body_text = driver.find_element(By.TAG_NAME, 'body').text.lower()
        if "thank you" in body_text or "thanks for subscribing" in body_text or "welcome" in body_text or "you're in" in body_text:
            update_task(task_id, message="   - ✅ Verified: Success message found.")
            return True
        else:
            update_task(task_id, message="   - ⚪ Submitted form, but could not verify success.")
            return False
    except:
        update_task(task_id, message="   - ❌ Could not find a newsletter signup form.")
        return False

def run_deep_simulation(url, task_id, num_products=3):
    update_task(task_id, status='running', message="--- Starting Full Spectrum Analysis ---", progress=5)
    
    driver = None
    try:
        options = webdriver.ChromeOptions()
        options.add_argument('--headless'); options.add_argument("--window-size=1920,1080"); options.add_argument('--no-sandbox'); options.add_argument('--disable-dev-shm-usage'); options.add_argument('user-agent=Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36')
        driver = webdriver.Chrome(options=options)
        driver.set_page_load_timeout(30)
        if not os.path.exists('screenshots'): os.makedirs('screenshots')

        capability_score = {'cdp_analytics': 0, 'ab_platform': 0, 'rec_widgets': 0}
        execution_score = {'homepage_change': 0, 'journey_personalization': 0}
        performance_score = {'fcp_score': 0, 'dom_interactive_score': 0}
        clarity_score = {'contrast_score': 0, 'color_score': 0}
        interactivity_score = {'add_to_cart': 0, 'newsletter_signup': 0}
        
        update_task(task_id, message="[Agent 1] Analyzing initial technology, performance, and design...", progress=10)
        driver.get(url)
        WebDriverWait(driver, 20).until(EC.presence_of_element_located((By.TAG_NAME, "body")))
        
        perf_data = get_performance_score(driver, task_id)
        performance_score['fcp_score'] = perf_data['fcp_score']
        performance_score['dom_interactive_score'] = perf_data['dom_interactive_score']
        perf_metrics = { "fcp": perf_data['fcp_value'], "dom_interactive": perf_data['dom_interactive_value'] }

        before_screenshot_path = f'screenshots/{task_id}_before.png'
        driver.save_screenshot(before_screenshot_path)
        
        clarity_data = get_visual_clarity_score(before_screenshot_path, task_id)
        clarity_score['contrast_score'] = clarity_data['contrast_score']
        clarity_score['color_score'] = clarity_data['color_score']
        
        found_vendors_list = detect_vendors(driver, task_id)
        for vendor in found_vendors_list:
            if "CDP" in vendor or "Analytics" in vendor: capability_score['cdp_analytics'] = 10
            if "A/B Testing" in vendor or "Personalization" in vendor: capability_score['ab_platform'] = 10
        if check_for_recommendations(driver, task_id): capability_score['rec_widgets'] = 5

        if attempt_newsletter_signup(driver, task_id):
            interactivity_score['newsletter_signup'] = 5
        
        update_task(task_id, message="[Agent 1] Simulating user journey...", progress=30)
        product_links = find_product_links(driver, url, task_id, limit=num_products)
        if not product_links: raise Exception("Could not find any product links to visit.")
        
        journey_recs_found = set()
        add_to_cart_success = False
        for i, link in enumerate(product_links, 1):
            driver.get(link)
            WebDriverWait(driver, 20).until(EC.presence_of_element_located((By.TAG_NAME, "body")))
            if check_for_recommendations(driver, task_id): journey_recs_found.add("product")
            if not add_to_cart_success:
                if attempt_add_to_cart(driver, task_id):
                    interactivity_score['add_to_cart'] = 10
                    add_to_cart_success = True
        update_task(task_id, progress=50)

        cart_paths = ['/cart', '/basket', '/checkout/cart']
        for path in cart_paths:
            cart_url = url.rstrip('/') + path
            try:
                driver.get(cart_url)
                WebDriverWait(driver, 10).until(EC.presence_of_element_located((By.TAG_NAME, "body")))
                if "cart" in driver.title.lower() or "basket" in driver.title.lower():
                    if check_for_recommendations(driver, task_id): journey_recs_found.add("cart")
                    break
            except Exception: continue
        update_task(task_id, progress=65)

        update_task(task_id, message="[Agent 1] Returning to homepage to check for dynamic changes...", progress=80)
        driver.get(url)
        WebDriverWait(driver, 20).until(EC.presence_of_element_located((By.TAG_NAME, "body")))
        after_screenshot_path = f'screenshots/{task_id}_after.png'
        driver.save_screenshot(after_screenshot_path)
        if check_for_recommendations(driver, task_id): journey_recs_found.add("homepage")

        difference = compare_images(before_screenshot_path, after_screenshot_path, f'screenshots/{task_id}_diff.png', task_id)
        execution_score['homepage_change'] = (difference / 5.0) * 20.0 if difference < 5.0 else 20.0
        execution_score['journey_personalization'] = (len(journey_recs_found) / 3) * 10
        update_task(task_id, message="[Agent 1] Quantitative analysis complete.", progress=90)
        
        total_capability_score = sum(capability_score.values())
        total_execution_score = sum(execution_score.values())
        total_performance_score = sum(performance_score.values())
        total_clarity_score = sum(clarity_score.values())
        total_interactivity_score = sum(interactivity_score.values())
        total_score = total_capability_score + total_execution_score + total_performance_score + total_clarity_score + total_interactivity_score

        score_report = {
            "url": url, "total_score": total_score, "capability_score": total_capability_score,
            "execution_score": total_execution_score, "performance_score": total_performance_score,
            "clarity_score": total_clarity_score, "interactivity_score": total_interactivity_score,
            "breakdown": {
                "capability": [
                    {"name": "CDP / Analytics Platform", "score": capability_score['cdp_analytics'], "max_score": 10},
                    {"name": "A/B & Personalization Platform", "score": capability_score['ab_platform'], "max_score": 10},
                    {"name": "Recommendation Widgets", "score": capability_score['rec_widgets'], "max_score": 5},
                ], "execution": [
                    {"name": "Dynamic Homepage Change", "score": execution_score['homepage_change'], "max_score": 20, "details": f"{difference:.2f}% visual change"},
                    {"name": "Journey Personalization", "score": execution_score['journey_personalization'], "max_score": 10, "details": f"{len(journey_recs_found)}/3 page types"},
                ], "performance": [
                    {"name": "First Contentful Paint (FCP)", "score": performance_score['fcp_score'], "max_score": 10, "details": f"{perf_metrics['fcp']:.0f}ms"},
                    {"name": "DOM Interactive Time", "score": performance_score['dom_interactive_score'], "max_score": 10, "details": f"{perf_metrics['dom_interactive']:.0f}ms"},
                ], "clarity": [
                    {"name": "Contrast Ratio", "score": clarity_score['contrast_score'], "max_score": 5, "details": f"{clarity_data['contrast_ratio']:.2f} ratio"},
                    {"name": "Color Palette Cohesion", "score": clarity_score['color_score'], "max_score": 10, "details": f"{clarity_data['num_colors']} dominant colors"},
                ], "interactivity": [
                    {"name": "Add to Cart", "score": interactivity_score['add_to_cart'], "max_score": 10},
                    {"name": "Newsletter Signup", "score": interactivity_score['newsletter_signup'], "max_score": 5},
                ]
            }
        }
        
        recommendations = generate_recommendations_with_llm(task_id, score_report)
        final_result = { "score_report": score_report, "recommendations": recommendations }
        
        update_task(task_id, message="[System] Generating PowerPoint report...", progress=98)
        presentation_filename = f"reports/Report_{task_id}.pptx"
        create_and_save_presentation(final_result, presentation_filename)
        update_task(task_id, message=f"[System] Report saved as {presentation_filename}")

        update_task(task_id, status='completed', message="--- All Agents Finished ---", progress=100, result=final_result)

    except Exception as e:
        update_task(task_id, status='failed', message=f"❌ An error occurred: {e}", progress=100)
    finally:
        if driver:
            driver.quit()

app = Flask(__name__)
CORS(app)

@app.route('/audit', methods=['POST'])
def start_audit():
    url = request.json.get('url')
    if not url: return jsonify({"error": "URL is required"}), 400
    task_id = str(uuid.uuid4())
    tasks[task_id] = {'status': 'pending', 'progress': 0, 'logs': [], 'result': None}
    thread = threading.Thread(target=run_deep_simulation, args=(url, task_id))
    thread.daemon = True
    thread.start()
    return jsonify({"task_id": task_id})

@app.route('/status/<task_id>', methods=['GET'])
def get_status(task_id):
    task = tasks.get(task_id)
    if not task: return jsonify({"error": "Task not found"}), 404
    return jsonify(task)

if __name__ == '__main__':
    initialize_knowledge_base()
    app.run(debug=True, port=5001)